"""
SRE Copilot — FastAPI service.

Endpoints:
  POST /alert           — Receives Alertmanager webhook payload (main entry point)
  POST /analyze         — Analyze a single alert (direct integration / testing)
  GET  /health          — Liveness probe
  GET  /status          — Ollama connectivity and model availability check
  POST /test-crash-loop — Simulate a CrashLoopBackOff alert (dev/demo)
  GET  /dashboard       — HTML dashboard with stats and incidents history
  GET  /api/incidents   — JSON list of all persisted investigations
  GET  /api/stats       — Aggregated stats (counts, MTTR, success rate)
  GET  /metrics         — Prometheus metrics (Mission 5)
"""
import asyncio
import json
import logging
import os
from collections import Counter
from contextlib import asynccontextmanager

import httpx
from fastapi import FastAPI
from fastapi.responses import HTMLResponse, Response, FileResponse, StreamingResponse
from io import BytesIO

from models import AlertManagerPayload, Alert, AgentResponse
from agent import run_sre_investigation
from tools.kubectl_tools import resolve_pod_name
import metrics

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
)
logger = logging.getLogger(__name__)

OLLAMA_BASE_URL  = os.getenv("OLLAMA_BASE_URL", "http://host.docker.internal:11434")
OLLAMA_MODEL     = os.getenv("OLLAMA_MODEL", "mistral")
GITOPS_REPO_PATH = os.getenv("GITOPS_REPO_PATH", "/gitops-repo")
INCIDENTS_LOG    = os.path.join(GITOPS_REPO_PATH, "incidents", "responses.jsonl")

# Semaphore: only 1 Mistral investigation at a time (CPU-bound on local host).
# Prevents memory exhaustion and model contention when Alertmanager sends bursts.
_investigation_lock = asyncio.Semaphore(1)


@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info(f"SRE Copilot starting — Ollama: {OLLAMA_BASE_URL} | Model: {OLLAMA_MODEL}")
    metrics.agent_info.labels(version="2.0.0", ollama_model=OLLAMA_MODEL).set(1)
    yield
    logger.info("SRE Copilot shutting down.")

app = FastAPI(
    title="SRE Copilot",
    description="AI-powered SRE agent for Kubernetes incident response and GitOps remediation",
    version="2.0.0",
    lifespan=lifespan,
)


# ── Health & Status ────────────────────────────────────────────────────────────

@app.get("/health")
async def health():
    return {"status": "ok", "service": "sre-copilot"}


@app.get("/status")
async def status():
    """Check Ollama connectivity and model availability."""
    try:
        async with httpx.AsyncClient(timeout=5) as client:
            resp = await client.get(f"{OLLAMA_BASE_URL}/api/tags")
        if resp.status_code == 200:
            models = [m["name"] for m in resp.json().get("models", [])]
            model_ok = any(OLLAMA_MODEL in m for m in models)
            return {
                "status":          "ok" if model_ok else "degraded",
                "ollama":          "reachable",
                "model":           OLLAMA_MODEL,
                "model_available": model_ok,
                "available_models": models,
            }
    except Exception as exc:
        return {
            "status": "degraded",
            "ollama": "unreachable",
            "error":  str(exc),
            "hint":   f"Run: ollama pull {OLLAMA_MODEL}",
        }


# ── Alert processing ───────────────────────────────────────────────────────────

@app.post("/alert", response_model=dict)
async def receive_alertmanager_webhook(payload: AlertManagerPayload):
    """
    Main webhook endpoint for Alertmanager.
    Processes all firing alerts sequentially under a semaphore so only one
    Mistral inference runs at a time (avoids CPU/memory contention).
    """
    firing = [a for a in payload.alerts if a.status == "firing"]
    if not firing:
        logger.info("Webhook received — no firing alerts, nothing to do.")
        return {"message": "No firing alerts.", "count": 0}

    logger.info(f"Webhook received — {len(firing)} firing alert(s)")
    results = []

    for alert in firing:
        logger.info(f"Queuing investigation: {alert.labels.alertname} / {alert.labels.namespace}")
        async with _investigation_lock:
            response = await run_sre_investigation(alert)
        results.append(response.model_dump())
        logger.info(
            f"Investigation complete: {alert.labels.alertname} → "
            f"classification={response.classification}"
        )

    return {
        "message": f"Processed {len(results)} alert(s)",
        "results": results,
    }


@app.post("/analyze", response_model=AgentResponse)
async def analyze_alert(alert: Alert):
    """Analyze a single alert directly. Used for testing and direct integrations."""
    logger.info(f"Direct analysis: {alert.labels.alertname}")
    async with _investigation_lock:
        return await run_sre_investigation(alert)


# ── Dev / Demo endpoints ───────────────────────────────────────────────────────

@app.post("/test-crash-loop")
async def test_crash_loop():
    """
    Simulate a CrashLoopBackOff alert using the real crash-app pod running in
    the demo namespace. Resolves the actual pod name dynamically so investigation
    data is real (not mocked).
    """
    pod_name = resolve_pod_name("demo", "crash-app")
    logger.info(f"test-crash-loop: resolved pod = {pod_name}")

    test_alert = Alert(
        status="firing",
        labels={
            "alertname": "PodRestartingTooMuch",
            "namespace":  "demo",
            "pod":         pod_name,
            "severity":   "critical",
        },
        annotations={
            "summary":     "Pod restarting frequently",
            "description": f"Pod {pod_name} in namespace demo is restarting too often.",
        },
    )
    async with _investigation_lock:
        return await run_sre_investigation(test_alert)


# ── Dashboard ──────────────────────────────────────────────────────────────────

def _read_incidents() -> list[dict]:
    """Read all persisted incidents from the JSONL log. Newest first."""
    if not os.path.exists(INCIDENTS_LOG):
        return []
    incidents = []
    with open(INCIDENTS_LOG, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            try:
                incidents.append(json.loads(line))
            except json.JSONDecodeError:
                continue
    incidents.sort(key=lambda x: x.get("timestamp", ""), reverse=True)
    return incidents


def _compute_stats(incidents: list[dict]) -> dict:
    """Aggregate stats across all incidents."""
    total = len(incidents)
    if total == 0:
        return {
            "total":              0,
            "auto_remediated":    0,
            "argocd_synced":      0,
            "success_rate":       0.0,
            "mttr_seconds":       0.0,
            "classifications":    {},
            "namespaces":         {},
            "severities":         {},
            "remediated_per_day": {},
        }

    auto_remediated = sum(1 for i in incidents if i.get("auto_remediation_applied"))
    argocd_synced   = sum(1 for i in incidents if i.get("argocd_sync_triggered"))
    durations       = [i.get("duration_seconds") for i in incidents if i.get("duration_seconds") is not None]
    mttr            = round(sum(durations) / len(durations), 2) if durations else 0.0

    classifications = Counter(i.get("classification", "Unknown") for i in incidents)
    namespaces      = Counter(i.get("namespace", "unknown")      for i in incidents)
    severities      = Counter(i.get("severity", "unknown")       for i in incidents)

    by_day = Counter()
    for i in incidents:
        ts = i.get("timestamp", "")
        if len(ts) >= 10:
            by_day[ts[:10]] += 1

    return {
        "total":              total,
        "auto_remediated":    auto_remediated,
        "argocd_synced":      argocd_synced,
        "success_rate":       round(100.0 * auto_remediated / total, 1) if total else 0.0,
        "mttr_seconds":       mttr,
        "classifications":    dict(classifications),
        "namespaces":         dict(namespaces),
        "severities":         dict(severities),
        "remediated_per_day": dict(sorted(by_day.items())),
    }


@app.get("/api/incidents")
async def list_incidents(limit: int = 100):
    """Return the last N incidents as JSON (newest first)."""
    incidents = _read_incidents()
    return {"count": len(incidents), "incidents": incidents[:limit]}


@app.get("/api/stats")
async def get_stats():
    """Return aggregated statistics across all incidents."""
    return _compute_stats(_read_incidents())


@app.get("/dashboard", response_class=HTMLResponse)
async def dashboard():
    """Serve the SRE Copilot dashboard HTML page."""
    return HTMLResponse(content=_DASHBOARD_HTML)


# ── Prometheus metrics (Mission 5) ─────────────────────────────────────────────

@app.get("/metrics")
async def prometheus_metrics():
    """Expose Prometheus-format metrics. Scraped by kube-prometheus-stack."""
    body, ctype = metrics.render_metrics()
    return Response(content=body, media_type=ctype)


# ── Static assets (PwC logo) ──────────────────────────────────────────────────

@app.get("/static/logo.png")
async def static_logo():
    """Sert le logo PwC officiel utilise par /dashboard et /audit-report."""
    return FileResponse("logo_pwc.png", media_type="image/png")


# ── Audit Report (PwC-style 24h alert report) ────────────────────────────────

from datetime import datetime, timedelta, timezone


def _filter_incidents_period(hours: int = 24) -> list[dict]:
    """Filtre les incidents sur les N dernières heures."""
    cutoff = datetime.now(timezone.utc) - timedelta(hours=hours)
    result = []
    for i in _read_incidents():
        ts_str = i.get("timestamp", "")
        if not ts_str:
            continue
        try:
            # Parse ISO format (avec ou sans timezone)
            ts = datetime.fromisoformat(ts_str.replace("Z", "+00:00"))
            if ts.tzinfo is None:
                ts = ts.replace(tzinfo=timezone.utc)
            if ts >= cutoff:
                result.append(i)
        except (ValueError, TypeError):
            continue
    return result


def _generate_audit_summary(incidents: list[dict], hours: int) -> dict:
    """Génère un résumé d'audit structuré façon PwC à partir des incidents."""
    total = len(incidents)
    if total == 0:
        return {
            "narrative": f"Aucune alerte n'a été reçue sur les {hours} dernières heures. "
                         f"Le système est en état nominal. Aucune intervention SRE requise.",
            "kpis": {"total": 0, "auto_remediated": 0, "success_rate": 0.0,
                     "mttr_seconds": 0.0, "argocd_synced": 0},
            "classifications": {},
            "severities":      {},
            "recommendations": [
                "Aucune action requise — surveillance active maintenue.",
            ],
            "risk_level": "low",
        }

    auto_remediated = sum(1 for i in incidents if i.get("auto_remediation_applied"))
    argocd_synced   = sum(1 for i in incidents if i.get("argocd_sync_triggered"))
    durations       = [i.get("duration_seconds") for i in incidents if i.get("duration_seconds") is not None]
    mttr            = round(sum(durations) / len(durations), 1) if durations else 0.0
    rate            = round(100.0 * auto_remediated / total, 1)

    classifications = Counter(i.get("classification", "Unknown") for i in incidents)
    severities      = Counter(i.get("severity", "info")          for i in incidents)
    top_class       = classifications.most_common(1)[0] if classifications else ("None", 0)
    critical_count  = severities.get("critical", 0)

    # Narrative en français pro
    narrative = (
        f"Pendant les {hours} dernières heures, l'agent SRE Copilot a traité "
        f"{total} incident(s). Le taux de remédiation automatique atteint {rate}% "
        f"({auto_remediated}/{total}). Le temps moyen de remédiation (MTTR) "
        f"s'élève à {mttr:.0f} secondes. La classification la plus fréquente "
        f"est « {top_class[0]} » ({top_class[1]} cas). "
        f"{critical_count} alerte(s) de sévérité critique ont été détectées."
    )

    # Recommandations générées à partir des patterns observés
    recommendations = []
    if rate < 80:
        recommendations.append(
            f"Le taux de remédiation automatique ({rate}%) est inférieur au seuil de 80%. "
            f"Identifier les classifications non-remédiables et étendre la couverture des fix_type."
        )
    if mttr > 200:
        recommendations.append(
            f"Le MTTR moyen ({mttr:.0f}s) dépasse 200 secondes, principalement à cause "
            f"du temps d'inférence Mistral 7B sur CPU. Envisager un GPU ou un modèle "
            f"plus léger pour la production."
        )
    if classifications.get("Unknown", 0) > total * 0.1:
        recommendations.append(
            f"{classifications['Unknown']} incident(s) classifiés Unknown. "
            f"Réviser le prompt LLM ou enrichir le keyword fallback."
        )
    if critical_count > total * 0.5:
        recommendations.append(
            f"Plus de 50% des alertes sont critiques. Analyser les causes racines récurrentes "
            f"et envisager des actions préventives en amont (Kyverno, quotas, limites)."
        )
    if argocd_synced == 0 and auto_remediated > 0:
        recommendations.append(
            "Aucune synchronisation ArgoCD n'a été déclenchée. Vérifier que l'agent "
            "a accès au token API ArgoCD et que l'Application demo-app est enregistrée."
        )
    if not recommendations:
        recommendations.append(
            "Le système opère dans les standards attendus. Maintenir la surveillance "
            "continue et conserver les garde-fous actuels (anti-boucle, validator)."
        )

    # Niveau de risque global (low / medium / high)
    if rate < 50 or critical_count > total * 0.7:
        risk = "high"
    elif rate < 80 or mttr > 250:
        risk = "medium"
    else:
        risk = "low"

    return {
        "narrative":       narrative,
        "kpis": {
            "total":           total,
            "auto_remediated": auto_remediated,
            "success_rate":    rate,
            "mttr_seconds":    mttr,
            "argocd_synced":   argocd_synced,
        },
        "classifications": dict(classifications),
        "severities":      dict(severities),
        "recommendations": recommendations,
        "risk_level":      risk,
    }


@app.get("/api/audit-report")
async def api_audit_report(hours: int = 24):
    """JSON audit report pour les N dernières heures (défaut 24h)."""
    incidents = _filter_incidents_period(hours)
    summary   = _generate_audit_summary(incidents, hours)
    now       = datetime.now(timezone.utc)
    return {
        "period_hours":  hours,
        "period_start":  (now - timedelta(hours=hours)).isoformat(),
        "period_end":    now.isoformat(),
        "generated_at":  now.isoformat(),
        "incidents":     incidents,
        "summary":       summary,
    }


@app.get("/audit-report", response_class=HTMLResponse)
async def audit_report_html():
    """Page HTML du rapport d'audit (imprimable Ctrl+P → PDF)."""
    return HTMLResponse(content=_AUDIT_REPORT_HTML)


# ── Export endpoints (/export/pptx  /export/pdf  /export/docx) ───────────────

@app.get("/export/pptx")
async def export_pptx(hours: int = 24):
    """Génère et télécharge un rapport PowerPoint."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    incidents = _filter_incidents_period(hours)
    summary   = _generate_audit_summary(incidents, hours)
    now       = datetime.now(timezone.utc)
    kpis      = summary["kpis"]

    ORANGE  = RGBColor(0xDC, 0x69, 0x00)
    ORANGED = RGBColor(0xC4, 0x4A, 0x00)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
    GREY    = RGBColor(0x66, 0x66, 0x66)
    LGREY   = RGBColor(0xF5, 0xF4, 0xF2)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_rect(sld, l, t, w, h, rgb):
        sp = sld.shapes.add_shape(1, l, t, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb
        sp.line.fill.background()

    def add_txt(sld, text, l, t, w, h, size=14, bold=False, rgb=BLACK, center=False, wrap=True):
        tb = sld.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb

    def tbl_cell(cell, text, rgb_bg=None, bold=False, size=12, rgb_fg=BLACK):
        cell.text = str(text)
        if rgb_bg:
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb_bg
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = rgb_fg

    # Slide 1 — Couverture (design PwC : fond pêche + formes oranges)
    s1 = prs.slides.add_slide(blank)
    # Fond crème blanc
    add_rect(s1, Inches(0), Inches(0), prs.slide_width, prs.slide_height, RGBColor(0xFF, 0xFC, 0xF8))
    # Dégradé pêche simulé coin haut-droit (3 rects transparents)
    add_rect(s1, Inches(8.0), Inches(0), Inches(5.33), Inches(4.0), RGBColor(0xFF, 0xE4, 0xC8))
    add_rect(s1, Inches(9.5), Inches(0), Inches(3.83), Inches(2.5), RGBColor(0xFF, 0xCE, 0xA8))
    add_rect(s1, Inches(11.0), Inches(0), Inches(2.33), Inches(1.5), RGBColor(0xFF, 0xB8, 0x88))
    # Remettre fond blanc par-dessus gauche pour adoucir
    add_rect(s1, Inches(0), Inches(0), Inches(8.5), Inches(5.0), RGBColor(0xFF, 0xFC, 0xF8))
    # Logo PwC (texte fallback car image difficile à embarquer en PPTX sans chemin)
    add_txt(s1, "pw", Inches(0.55), Inches(0.3), Inches(0.7), Inches(0.6), size=28, bold=True, rgb=BLACK)
    add_txt(s1, "c",  Inches(1.15), Inches(0.3), Inches(0.4), Inches(0.6), size=28, bold=True, rgb=ORANGE)
    # Date haut-droite
    add_txt(s1, now.strftime("%d/%m/%Y %H:%M"), Inches(9.0), Inches(0.35), Inches(4.0), Inches(0.5), size=10, rgb=GREY)
    # Grand titre
    add_txt(s1, "Rapport", Inches(0.55), Inches(1.5), Inches(8.0), Inches(1.2), size=58, bold=True, rgb=BLACK)
    add_txt(s1, "D'audit", Inches(0.55), Inches(2.6), Inches(8.0), Inches(1.2), size=58, bold=True, rgb=BLACK)
    # Sous-titre
    add_txt(s1, "Alertes Kubernetes & remédiation SRE", Inches(0.55), Inches(4.0), Inches(9.0), Inches(0.6), size=16, bold=True, rgb=RGBColor(0x28,0x28,0x28))
    # Métadonnées
    add_txt(s1, f"Période : {hours}h  ·  Généré le {now.strftime('%d/%m/%Y à %H:%M')} UTC",
            Inches(0.55), Inches(4.55), Inches(9.0), Inches(0.5), size=11, rgb=GREY)
    # Forme gauche (orange foncé) — parallélogramme bas
    from pptx.util import Emu
    sp_back = s1.shapes.add_shape(6, Inches(0), Inches(5.35), Inches(7.5), Inches(0.88))
    sp_back.fill.solid(); sp_back.fill.fore_color.rgb = ORANGED
    sp_back.line.fill.background()
    # Forme droite (orange vif)
    sp_front = s1.shapes.add_shape(6, Inches(2.8), Inches(5.0), Inches(7.5), Inches(0.88))
    sp_front.fill.solid(); sp_front.fill.fore_color.rgb = ORANGE
    sp_front.line.fill.background()

    # Slide 2 — KPIs
    s2 = prs.slides.add_slide(blank)
    add_rect(s2, Inches(0), Inches(0), prs.slide_width, Inches(1.1), ORANGE)
    add_txt(s2, "Résumé exécutif", Inches(0.5), Inches(0.15), Inches(9), Inches(0.8), size=28, bold=True, rgb=WHITE)
    for i, (lbl, val) in enumerate([
        ("Incidents traités",  kpis["total"]),
        ("Auto-remédiés",      kpis["auto_remediated"]),
        ("Taux de succès",     f"{kpis['success_rate']}%"),
        ("MTTR moyen",         f"{kpis['mttr_seconds']:.0f}s"),
    ]):
        x = Inches(0.4 + i * 3.2)
        add_rect(s2, x, Inches(1.3), Inches(3.0), Inches(1.9), LGREY)
        add_txt(s2, str(val), x, Inches(1.45), Inches(3.0), Inches(1.0), size=38, bold=True, rgb=ORANGE, center=True)
        add_txt(s2, lbl,      x, Inches(2.55), Inches(3.0), Inches(0.5), size=12, rgb=GREY, center=True)
    add_txt(s2, summary["narrative"], Inches(0.5), Inches(3.4), Inches(12.3), Inches(2.5), size=13, wrap=True)

    # Slide 3 — Classifications
    s3 = prs.slides.add_slide(blank)
    add_rect(s3, Inches(0), Inches(0), prs.slide_width, Inches(1.1), ORANGE)
    add_txt(s3, "Répartition par classification", Inches(0.5), Inches(0.15), Inches(10), Inches(0.8), size=28, bold=True, rgb=WHITE)
    classes = sorted(summary["classifications"].items(), key=lambda x: -x[1])
    if classes:
        rows = len(classes) + 1
        tbl = s3.shapes.add_table(rows, 2, Inches(1.0), Inches(1.4), Inches(7.0),
                                   Inches(min(5.5, rows * 0.55))).table
        tbl.columns[0].width = Inches(5); tbl.columns[1].width = Inches(2)
        tbl_cell(tbl.cell(0, 0), "Classification", ORANGE, True, 13, WHITE)
        tbl_cell(tbl.cell(0, 1), "Incidents",      ORANGE, True, 13, WHITE)
        for ri, (cls, cnt) in enumerate(classes, 1):
            bg = LGREY if ri % 2 == 0 else None
            tbl_cell(tbl.cell(ri, 0), cls,      bg, False, 12)
            tbl_cell(tbl.cell(ri, 1), str(cnt), bg, False, 12)

    # Slide 4 — Journal des incidents
    s4 = prs.slides.add_slide(blank)
    add_rect(s4, Inches(0), Inches(0), prs.slide_width, Inches(1.1), ORANGE)
    add_txt(s4, "Journal des incidents", Inches(0.5), Inches(0.15), Inches(10), Inches(0.8), size=28, bold=True, rgb=WHITE)
    recent = incidents[:15]
    if recent:
        rows = len(recent) + 1
        tbl = s4.shapes.add_table(rows, 4, Inches(0.2), Inches(1.3), Inches(12.9),
                                   Inches(min(5.9, rows * 0.42))).table
        for c, w in enumerate([Inches(2.2), Inches(3.0), Inches(2.0), Inches(5.7)]):
            tbl.columns[c].width = w
        for c, h in enumerate(["Horodatage", "Classification", "Sévérité", "Pod"]):
            tbl_cell(tbl.cell(0, c), h, ORANGE, True, 12, WHITE)
        for ri, inc in enumerate(recent, 1):
            bg = LGREY if ri % 2 == 0 else None
            for c, val in enumerate([
                inc.get("timestamp", "")[:16].replace("T", " "),
                inc.get("classification", "Unknown"),
                inc.get("severity", "—"),
                inc.get("pod_name", inc.get("pod", "—")),
            ]):
                tbl_cell(tbl.cell(ri, c), val, bg, False, 11)

    buf = BytesIO(); prs.save(buf); buf.seek(0)
    fname = f"SRE-Copilot-{hours}h-{now.strftime('%Y%m%d-%H%M')}.pptx"
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'})


@app.get("/export/pdf")
async def export_pdf(hours: int = 24):
    """Génère et télécharge un rapport PDF avec design PwC fidèle à la page HTML."""
    from fpdf import FPDF

    incidents = _filter_incidents_period(hours)
    summary   = _generate_audit_summary(incidents, hours)
    now       = datetime.now(timezone.utc)
    kpis      = summary["kpis"]

    FONT_DIR = "/usr/share/fonts/truetype/dejavu"

    class AuditPDF(FPDF):
        def footer(self):
            self.set_y(-12)
            self.set_font("Body", "", 7.5)
            self.set_text_color(160, 160, 160)
            self.cell(0, 8, f"SRE Copilot  ·  Rapport d'audit {hours}h  ·  Page {self.page_no()}", align="C")

    pdf = AuditPDF()
    pdf.add_font("Body",  "",  f"{FONT_DIR}/DejaVuSans.ttf")
    pdf.add_font("Body",  "B", f"{FONT_DIR}/DejaVuSans-Bold.ttf")
    pdf.add_font("Body",  "I", f"{FONT_DIR}/DejaVuSans-Oblique.ttf")
    pdf.add_font("Serif", "",  f"{FONT_DIR}/DejaVuSerif.ttf")
    pdf.add_font("Serif", "B", f"{FONT_DIR}/DejaVuSerif-Bold.ttf")
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.set_margins(15, 15, 15)

    # ── Palette exacte du HTML ──
    ORANGE   = (220, 105,   0)   # --orange principal
    ORANGE_D = (180,  64,   0)   # forme gauche (foncée)
    ORANGE_B = (232,  93,   4)   # forme droite (vive)
    CREAM_R  = (255, 248, 244)   # fond pêche coin haut-droit (simulé par un dégradé de rects)
    CREAM_W  = (255, 253, 251)   # fond crème blanc
    GREY     = (100, 100, 100)
    INK      = ( 26,  26,  26)
    ROW_ALT  = (248, 245, 242)
    BORDER   = (220, 220, 220)

    # ── Helper: bande orange en-tête de page de contenu ──
    def page_header(title):
        # Bande orange pleine largeur
        pdf.set_fill_color(*ORANGE)
        pdf.rect(0, 0, 210, 18, "F")
        pdf.set_xy(14, 3.5)
        pdf.set_font("Serif", "B", 15)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(0, 11, title)
        # Breadcrumb discret
        pdf.set_xy(14, 22)
        pdf.set_font("Body", "", 7.5)
        pdf.set_text_color(*GREY)
        pdf.cell(0, 5, "RAPPORT D'AUDIT  ·  ALERTES KUBERNETES & REMÉDIATION SRE")
        pdf.set_text_color(*INK)
        pdf.set_xy(14, 32)

    # ══════════════════════════════════════════════════════
    # PAGE 1 — COUVERTURE (fidèle au PDF du design HTML)
    # ══════════════════════════════════════════════════════
    pdf.add_page()

    # Fond blanc crème sur toute la page
    pdf.set_fill_color(*CREAM_W)
    pdf.rect(0, 0, 210, 297, "F")

    # Dégradé pêche simulé coin haut-droit (3 rectangles de plus en plus transparents)
    pdf.set_fill_color(255, 228, 200)
    pdf.rect(140, 0, 70, 80, "F")
    pdf.set_fill_color(255, 238, 218)
    pdf.rect(120, 0, 90, 60, "F")
    pdf.set_fill_color(255, 248, 240)
    pdf.rect(100, 0, 110, 40, "F")
    # Remettre le fond crème par-dessus pour adoucir
    pdf.set_fill_color(255, 252, 248)
    pdf.rect(0, 0, 105, 100, "F")

    # ── Logo PwC haut-gauche ──
    try:
        pdf.image("logo_pwc.png", 14, 16, w=26)
    except Exception:
        # Fallback texte si logo absent
        pdf.set_font("Serif", "B", 20)
        pdf.set_text_color(40, 40, 40)
        pdf.set_xy(14, 14)
        pdf.cell(15, 10, "pw")
        pdf.set_text_color(*ORANGE_B)
        pdf.set_xy(29, 14)
        pdf.cell(8, 10, "c")

    # ── Date haut-droite ──
    pdf.set_font("Body", "", 9)
    pdf.set_text_color(*GREY)
    pdf.set_xy(110, 18)
    pdf.cell(86, 7, now.strftime("%d/%m/%Y %H:%M"), align="R")

    # ── Grand titre Serif (exact comme le HTML .cover-h1) ──
    pdf.set_text_color(*INK)
    pdf.set_font("Serif", "B", 58)
    pdf.set_xy(14, 68)
    pdf.cell(0, 26, "Rapport")
    pdf.set_xy(14, 94)
    pdf.cell(0, 26, "D'audit")

    # ── Sous-titre en gras ──
    pdf.set_font("Body", "B", 14)
    pdf.set_text_color(40, 40, 40)
    pdf.set_xy(14, 130)
    pdf.cell(0, 8, "Alertes Kubernetes & remédiation SRE")

    # ── Métadonnées ──
    pdf.set_font("Body", "", 9)
    pdf.set_text_color(*GREY)
    pdf.set_xy(14, 143)
    pdf.cell(90, 6, f"Période : {hours}h")
    pdf.set_xy(14, 150)
    pdf.cell(90, 6, f"Généré le {now.strftime('%d/%m/%Y à %H:%M')} UTC")

    # ── Deux parallélogrammes PwC bas de page ──
    # Position: ~57% de 297mm ≈ y=169, deux formes qui se chevauchent
    # Forme gauche (foncée, déborde à gauche)
    pdf.set_fill_color(*ORANGE_D)
    pdf.polygon([(-8, 192), (148, 192), (134, 218), (-22, 218)], fill=True)
    # Forme droite (orange vive, devant)
    pdf.set_fill_color(*ORANGE_B)
    pdf.polygon([(88, 178), (244, 178), (230, 204), ( 74, 204)], fill=True)

    # ══════════════════════════════════════════════════════
    # PAGE 2 — RÉSUMÉ EXÉCUTIF
    # ══════════════════════════════════════════════════════
    pdf.add_page()
    page_header("Résumé exécutif")

    # 4 KPI cards côte à côte (comme le HTML .kpi-grid)
    kpi_colors = [
        (30,  96, 128),   # bleu acier   — total
        (139, 31,  31),   # bordeaux     — auto-remédiés
        (196, 74,   0),   # orange foncé — taux
        (232, 93,   4),   # orange vif   — MTTR
    ]
    kpi_data = [
        ("Total incidents",  str(kpis["total"])),
        ("Auto-remédiés",    str(kpis["auto_remediated"])),
        ("Taux de succès",   f"{kpis['success_rate']}%"),
        ("MTTR moyen",       f"{kpis['mttr_seconds']:.0f}s"),
    ]
    card_w, card_h = 43, 30
    for i, ((lbl, val), color) in enumerate(zip(kpi_data, kpi_colors)):
        bx = 14 + i * (card_w + 2)
        # Fond coloré de la carte
        pdf.set_fill_color(*color)
        pdf.rect(bx, 32, card_w, card_h, "F")
        # Valeur grande
        pdf.set_xy(bx, 34)
        pdf.set_font("Serif", "B", 22)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(card_w, 14, val, align="C")
        # Label
        pdf.set_xy(bx, 50)
        pdf.set_font("Body", "", 7.5)
        pdf.set_text_color(255, 255, 220)
        pdf.cell(card_w, 7, lbl, align="C")

    # Narrative
    pdf.set_xy(14, 72)
    pdf.set_font("Body", "", 10.5)
    pdf.set_text_color(*INK)
    pdf.multi_cell(182, 6, summary["narrative"])

    # Niveau de risque (bandeau coloré)
    risk = summary.get("risk_level", "low")
    risk_colors = {"low": (45, 122, 62), "medium": (184, 110, 0), "high": (184, 27, 27)}
    risk_bgs    = {"low": (238, 247, 239), "medium": (255, 244, 229), "high": (251, 240, 240)}
    risk_labels = {"low": "Niveau de risque · FAIBLE", "medium": "Niveau de risque · MODÉRÉ", "high": "Niveau de risque · ÉLEVÉ"}
    rc  = risk_colors.get(risk, (45, 122, 62))
    rbg = risk_bgs.get(risk, (238, 247, 239))
    rl  = risk_labels.get(risk, "—")
    cy  = pdf.get_y() + 6
    pdf.set_fill_color(*rbg)
    pdf.rect(14, cy, 182, 12, "F")
    pdf.set_fill_color(*rc)
    pdf.rect(14, cy, 3, 12, "F")
    pdf.set_xy(20, cy + 2)
    pdf.set_font("Body", "B", 9)
    pdf.set_text_color(*rc)
    pdf.cell(0, 7, rl)

    # ══════════════════════════════════════════════════════
    # PAGE 3 — CLASSIFICATIONS
    # ══════════════════════════════════════════════════════
    pdf.add_page()
    page_header("Répartition par classification")
    classes = sorted(summary["classifications"].items(), key=lambda x: -x[1])
    if classes:
        # En-tête tableau
        pdf.set_font("Body", "B", 10)
        pdf.set_fill_color(*ORANGE)
        pdf.set_text_color(255, 255, 255)
        pdf.cell(140, 10, "  Classification", border=0, fill=True)
        pdf.cell( 40, 10, "Incidents", border=0, fill=True, align="C")
        pdf.ln()
        pdf.set_font("Body", "", 10)
        for i, (cls, cnt) in enumerate(classes):
            bg = ROW_ALT if i % 2 == 0 else (255, 255, 255)
            pdf.set_fill_color(*bg)
            pdf.set_text_color(*INK)
            pdf.cell(140, 9, "  " + cls, border=0, fill=True)
            pdf.cell( 40, 9, str(cnt), border=0, fill=True, align="C")
            pdf.ln()
    else:
        pdf.set_font("Body", "I", 10)
        pdf.set_text_color(*GREY)
        pdf.cell(0, 8, "Aucune classification sur la période.")

    # ══════════════════════════════════════════════════════
    # PAGE 4 — JOURNAL DES INCIDENTS
    # ══════════════════════════════════════════════════════
    pdf.add_page()
    page_header("Journal des incidents")
    recent = incidents[:30]
    if recent:
        cols = [
            ("Horodatage",     35),
            ("Classification", 48),
            ("Sévérité",       26),
            ("Pod",            70),
            ("Rem.",           11),
        ]
        pdf.set_font("Body", "B", 8.5)
        pdf.set_fill_color(*ORANGE)
        pdf.set_text_color(255, 255, 255)
        for h, w in cols:
            pdf.cell(w, 9, " " + h, border=0, fill=True)
        pdf.ln()
        pdf.set_font("Body", "", 8)
        for i, inc in enumerate(recent):
            pdf.set_fill_color(*(ROW_ALT if i % 2 == 0 else (255, 255, 255)))
            pdf.set_text_color(*INK)
            pod = str(inc.get("pod_name", inc.get("pod", "—")))
            rem = "Oui" if inc.get("auto_remediation_applied") else "Non"
            for val, w in [
                (inc.get("timestamp", "")[:16].replace("T", " "), 35),
                (str(inc.get("classification", "Unknown"))[:26],   48),
                (str(inc.get("severity", "—")),                    26),
                (pod[:36],                                          70),
                (rem,                                               11),
            ]:
                pdf.cell(w, 7, " " + str(val), border=0, fill=True)
            pdf.ln()
    else:
        pdf.set_font("Body", "I", 10)
        pdf.set_text_color(*GREY)
        pdf.cell(0, 8, "Aucun incident enregistré sur cette période.")

    buf = BytesIO(pdf.output())
    fname = f"SRE-Copilot-{hours}h-{now.strftime('%Y%m%d-%H%M')}.pdf"
    return StreamingResponse(buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'})


@app.get("/export/docx")
async def export_docx(hours: int = 24):
    """Génère et télécharge un rapport Word."""
    from docx import Document
    from docx.shared import Pt, RGBColor as DocRGB, Cm

    incidents = _filter_incidents_period(hours)
    summary   = _generate_audit_summary(incidents, hours)
    now       = datetime.now(timezone.utc)
    kpis      = summary["kpis"]

    ORANGE  = DocRGB(0xDC, 0x69, 0x00)
    ORANGED = DocRGB(0xB4, 0x40, 0x00)
    BLACK   = DocRGB(0x1A, 0x1A, 0x1A)
    GREY    = DocRGB(0x64, 0x64, 0x64)

    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Cm(2.0)
        sec.bottom_margin = Cm(2.0)
        sec.left_margin   = sec.right_margin = Cm(2.5)

    def set_para_spacing(p, before=0, after=0):
        pPr = p._p.get_or_add_pPr()
        spacing = OxmlElement('w:spacing')
        spacing.set(qn('w:before'), str(before))
        spacing.set(qn('w:after'), str(after))
        pPr.append(spacing)

    def set_para_shading(p, fill_hex):
        pPr = p._p.get_or_add_pPr()
        shd = OxmlElement('w:shd')
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), fill_hex)
        pPr.append(shd)

    # ── Logo PwC (texte stylisé) ──
    p_logo = doc.add_paragraph()
    set_para_spacing(p_logo, before=0, after=60)
    r_pw = p_logo.add_run("pw")
    r_pw.font.size = Pt(28); r_pw.font.bold = True; r_pw.font.color.rgb = BLACK
    r_c = p_logo.add_run("c")
    r_c.font.size = Pt(28); r_c.font.bold = True; r_c.font.color.rgb = ORANGE

    # ── Ligne de séparation orange ──
    p_line = doc.add_paragraph()
    set_para_spacing(p_line, before=0, after=280)
    p_line.paragraph_format.border_bottom = True
    pPr = p_line._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '12')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), 'DC6900')
    pBdr.append(bottom)
    pPr.append(pBdr)

    # ── Grand titre ──
    p_t1 = doc.add_paragraph()
    set_para_spacing(p_t1, before=0, after=0)
    r_t1 = p_t1.add_run("Rapport")
    r_t1.font.size = Pt(48); r_t1.font.bold = True; r_t1.font.color.rgb = BLACK
    r_t1.font.name = "Georgia"

    p_t2 = doc.add_paragraph()
    set_para_spacing(p_t2, before=0, after=200)
    r_t2 = p_t2.add_run("D'audit")
    r_t2.font.size = Pt(48); r_t2.font.bold = True; r_t2.font.color.rgb = BLACK
    r_t2.font.name = "Georgia"

    # ── Sous-titre ──
    p_sub = doc.add_paragraph()
    set_para_spacing(p_sub, before=0, after=100)
    r_sub = p_sub.add_run("Alertes Kubernetes & remédiation SRE")
    r_sub.font.size = Pt(14); r_sub.font.bold = True; r_sub.font.color.rgb = BLACK

    # ── Métadonnées ──
    p_meta = doc.add_paragraph()
    set_para_spacing(p_meta, before=0, after=400)
    r_meta = p_meta.add_run(
        f"Période : {hours} dernières heures   ·   "
        f"Généré le {now.strftime('%d/%m/%Y à %H:%M')} UTC   ·   Confidentiel"
    )
    r_meta.font.size = Pt(9); r_meta.font.color.rgb = GREY

    # ── Bloc orange décoratif (simule les formes PwC) ──
    p_bar1 = doc.add_paragraph()
    set_para_spacing(p_bar1, before=0, after=8)
    set_para_shading(p_bar1, 'B44000')
    r_bar1 = p_bar1.add_run("  ")
    r_bar1.font.size = Pt(16)

    p_bar2 = doc.add_paragraph()
    set_para_spacing(p_bar2, before=0, after=600)
    set_para_shading(p_bar2, 'E85D04')
    r_bar2 = p_bar2.add_run("  ")
    r_bar2.font.size = Pt(16)

    doc.add_page_break()

    # Section 1 — Résumé
    h1 = doc.add_heading("1. Résumé exécutif", level=1)
    h1.runs[0].font.color.rgb = ORANGE
    doc.add_paragraph(summary["narrative"])
    doc.add_paragraph()
    kpi_tbl = doc.add_table(2, 4); kpi_tbl.style = "Table Grid"
    for i, (lbl, val) in enumerate([
        ("Incidents traités",  str(kpis["total"])),
        ("Auto-remédiés",      str(kpis["auto_remediated"])),
        ("Taux de succès",     f"{kpis['success_rate']}%"),
        ("MTTR moyen",         f"{kpis['mttr_seconds']:.0f}s"),
    ]):
        kpi_tbl.cell(0, i).text = lbl
        kpi_tbl.cell(0, i).paragraphs[0].runs[0].font.bold = True
        kpi_tbl.cell(1, i).text = val
        for run in kpi_tbl.cell(1, i).paragraphs[0].runs:
            run.font.size = Pt(16); run.font.bold = True; run.font.color.rgb = ORANGE
    doc.add_paragraph()

    # Section 2 — Classifications
    h2 = doc.add_heading("2. Répartition par classification", level=1)
    h2.runs[0].font.color.rgb = ORANGE
    classes = sorted(summary["classifications"].items(), key=lambda x: -x[1])
    if classes:
        cls_tbl = doc.add_table(len(classes) + 1, 2); cls_tbl.style = "Table Grid"
        for c, h in enumerate(["Classification", "Incidents"]):
            cls_tbl.cell(0, c).text = h
            cls_tbl.cell(0, c).paragraphs[0].runs[0].font.bold = True
        for ri, (cls, cnt) in enumerate(classes, 1):
            cls_tbl.cell(ri, 0).text = cls; cls_tbl.cell(ri, 1).text = str(cnt)
    doc.add_paragraph()

    # Section 3 — Incidents
    h3 = doc.add_heading("3. Journal détaillé des incidents", level=1)
    h3.runs[0].font.color.rgb = ORANGE
    sample = incidents[:50]
    if sample:
        inc_tbl = doc.add_table(len(sample) + 1, 5); inc_tbl.style = "Table Grid"
        for c, h in enumerate(["Horodatage", "Classification", "Sévérité", "Pod", "Remédiation"]):
            inc_tbl.cell(0, c).text = h
            inc_tbl.cell(0, c).paragraphs[0].runs[0].font.bold = True
        for ri, inc in enumerate(sample, 1):
            for c, val in enumerate([
                inc.get("timestamp", "")[:16].replace("T", " "),
                inc.get("classification", "Unknown"),
                inc.get("severity", "—"),
                inc.get("pod_name", inc.get("pod", "—")),
                "Oui" if inc.get("auto_remediation_applied") else "Non",
            ]):
                inc_tbl.cell(ri, c).text = str(val)
                for run in inc_tbl.cell(ri, c).paragraphs[0].runs:
                    run.font.size = Pt(9)
    else:
        doc.add_paragraph("Aucun incident enregistré sur cette période.")

    buf = BytesIO(); doc.save(buf); buf.seek(0)
    fname = f"SRE-Copilot-{hours}h-{now.strftime('%Y%m%d-%H%M')}.docx"
    return StreamingResponse(buf,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'})


# ── HTML pages servies par /dashboard et /audit-report ────────────────────
# Externalisees dans html_pages.py pour separer code Python (lisible) et HTML (volumineux)
from html_pages import DASHBOARD_HTML as _DASHBOARD_HTML, AUDIT_REPORT_HTML as _AUDIT_REPORT_HTML
